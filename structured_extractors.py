"""
结构化文档提取模块

支持从矢量PDF、Word、PPTX中提取结构化内容并转换为Markdown
"""

import re
from abc import ABC, abstractmethod
from pathlib import Path
from typing import List, Dict, Any, Optional, Tuple
from dataclasses import dataclass
from datetime import datetime


@dataclass
class ExtractedContent:
    """提取的内容结果"""
    title: str
    content: str
    metadata: Dict[str, Any]
    source_type: str  # 'vector_pdf', 'word', 'pptx'


class BaseExtractor(ABC):
    """文档提取器基类"""

    @abstractmethod
    def extract(self, file_path: Path) -> ExtractedContent:
        """提取文档内容"""
        pass

    def _create_header(self, file_path: Path, source_type: str, metadata: Dict[str, Any] = None) -> str:
        """创建统一的文档头部"""
        title = file_path.stem

        lines = [
            f"# {title}",
            "",
            "> **转换信息**",
            f"> - 源文件：{file_path.name}",
            f"> - 文档类型：{source_type}",
            f"> - 处理时间：{datetime.now().strftime('%Y-%m-%d %H:%M')}",
        ]

        if metadata:
            for key, value in metadata.items():
                if value:
                    lines.append(f"> - {key}：{value}")

        lines.extend(["", "---", ""])
        return "\n".join(lines)

    def _escape_table_cell(self, text: str) -> str:
        """转义表格单元格中的特殊字符"""
        return text.replace('|', '\\|').replace('\n', '<br>')


class VectorPDFExtractor(BaseExtractor):
    """矢量PDF提取器 - 使用PyMuPDF提取结构化内容"""

    # 标题大小阈值（可根据需要调整）
    H1_MIN_SIZE = 18
    H2_MIN_SIZE = 14
    H3_MIN_SIZE = 12

    def __init__(self):
        self._fitz = None

    def _get_fitz(self):
        """延迟加载 fitz 模块"""
        if self._fitz is None:
            try:
                import fitz
                self._fitz = fitz
            except ImportError:
                raise ImportError("pymupdf 未安装，请运行: pip install pymupdf")
        return self._fitz

    def extract(self, file_path: Path) -> ExtractedContent:
        """提取矢量PDF的结构化内容"""
        fitz = self._get_fitz()
        doc = fitz.open(str(file_path))

        try:
            all_blocks = []
            tables_found = []
            images_found = []

            for page_num, page in enumerate(doc, 1):
                # 提取带格式的文本块
                text_dict = page.get_text("dict")

                page_blocks = []
                for block in text_dict.get("blocks", []):
                    if block.get("type") == 0:  # 文本块
                        block_info = self._parse_text_block(block, page_num)
                        if block_info:
                            page_blocks.append(block_info)
                    elif block.get("type") == 1:  # 图片块
                        image_info = self._parse_image_block(block, page_num)
                        if image_info:
                            images_found.append(image_info)

                # 检测表格
                tables = page.find_tables()
                if tables:
                    for table in tables.tables:
                        table_md = self._convert_table_to_markdown(table)
                        tables_found.append({
                            'page': page_num,
                            'markdown': table_md,
                            'bbox': table.bbox
                        })

                all_blocks.extend(sorted(page_blocks, key=lambda x: (x['y0'], x['x0'])))

            # 构建Markdown内容
            content_parts = []
            current_page = 0

            for block in all_blocks:
                if block['page'] != current_page:
                    current_page = block['page']
                    content_parts.append(f"\n<!-- 第{current_page}页 -->\n")

                # 检查此文本块附近是否有图片
                nearby_images = self._find_nearby_images(block, images_found)
                for img in nearby_images:
                    content_parts.append(self._image_to_placeholder(img))

                md_text = self._block_to_markdown(block)
                if md_text:
                    content_parts.append(md_text)

            # 合并内容
            content = "\n\n".join(content_parts)
            content = self._clean_content(content)

            # 元数据
            metadata = {
                '总页数': len(doc),
                '提取字符数': len(content),
            }

            header = self._create_header(file_path, "矢量PDF", metadata)
            full_content = header + "\n\n" + content

            return ExtractedContent(
                title=file_path.stem,
                content=full_content,
                metadata=metadata,
                source_type='vector_pdf'
            )

        finally:
            doc.close()

    def _parse_text_block(self, block: Dict, page_num: int) -> Optional[Dict]:
        """解析文本块"""
        lines = block.get("lines", [])
        if not lines:
            return None

        texts = []
        max_size = 0
        is_bold = False

        for line in lines:
            line_text = ""
            for span in line.get("spans", []):
                text = span.get("text", "").strip()
                if text:
                    line_text += text + " "
                    size = span.get("size", 11)
                    flags = span.get("flags", 0)
                    if size > max_size:
                        max_size = size
                    if flags & 2 ** 4:  # 粗体标志
                        is_bold = True

            if line_text.strip():
                texts.append(line_text.strip())

        if not texts:
            return None

        full_text = " ".join(texts)

        return {
            'text': full_text,
            'size': max_size,
            'is_bold': is_bold,
            'x0': block.get('bbox', [0, 0, 0, 0])[0],
            'y0': block.get('bbox', [0, 0, 0, 0])[1],
            'page': page_num,
            'type': self._determine_block_type(full_text, max_size, is_bold)
        }

    def _determine_block_type(self, text: str, size: float, is_bold: bool) -> str:
        """确定文本块类型"""
        # 检查是否为列表项
        if re.match(r'^[\-\*•・]\s+', text):
            return 'bullet_list'
        if re.match(r'^\d+[\.\)]\s+', text):
            return 'numbered_list'

        # 根据字体大小判断标题级别
        if size >= self.H1_MIN_SIZE or (is_bold and size >= self.H2_MIN_SIZE):
            return 'h1'
        elif size >= self.H2_MIN_SIZE:
            return 'h2'
        elif size >= self.H3_MIN_SIZE:
            return 'h3'

        return 'paragraph'

    def _block_to_markdown(self, block: Dict) -> str:
        """将文本块转换为Markdown"""
        text = block['text']
        block_type = block['type']

        if not text.strip():
            return ""

        if block_type == 'h1':
            return f"# {text}"
        elif block_type == 'h2':
            return f"## {text}"
        elif block_type == 'h3':
            return f"### {text}"
        elif block_type == 'bullet_list':
            # 移除原有的列表标记，添加Markdown格式
            clean_text = re.sub(r'^[\-\*•・]\s+', '', text)
            return f"- {clean_text}"
        elif block_type == 'numbered_list':
            # 保留数字
            return text
        else:
            return text

    def _convert_table_to_markdown(self, table) -> str:
        """将PyMuPDF表格转换为Markdown"""
        try:
            df = table.to_pandas()
            return self._dataframe_to_markdown(df)
        except:
            # 降级处理：直接获取文本
            rows = []
            for row in table.extract():
                if row:
                    cells = [self._escape_table_cell(str(cell or "")) for cell in row]
                    rows.append("| " + " | ".join(cells) + " |")

            if not rows:
                return ""

            # 添加表头分隔行
            if len(rows) > 0:
                col_count = rows[0].count("|") - 1
                separator = "|" + "---|" * col_count
                rows.insert(1, separator)

            return "\n".join(rows)

    def _dataframe_to_markdown(self, df) -> str:
        """将DataFrame转换为Markdown表格"""
        lines = []

        # 表头
        header = "| " + " | ".join(self._escape_table_cell(str(col)) for col in df.columns) + " |"
        lines.append(header)

        # 分隔行
        separator = "|" + "---|" * len(df.columns)
        lines.append(separator)

        # 数据行
        for _, row in df.iterrows():
            row_str = "| " + " | ".join(self._escape_table_cell(str(cell)) for cell in row) + " |"
            lines.append(row_str)

        return "\n".join(lines)

    def _parse_image_block(self, block: Dict, page_num: int) -> Optional[Dict]:
        """解析图片块"""
        bbox = block.get('bbox', [0, 0, 0, 0])
        return {
            'page': page_num,
            'x0': bbox[0],
            'y0': bbox[1],
            'x1': bbox[2],
            'y1': bbox[3],
            'type': 'image'
        }

    def _find_nearby_images(self, text_block: Dict, images: List[Dict], threshold: float = 50) -> List[Dict]:
        """查找文本块附近的图片（避免重复输出）"""
        nearby = []
        text_y = text_block.get('y0', 0)
        text_page = text_block.get('page', 0)

        for img in images:
            if img.get('page') != text_page:
                continue
            # 如果图片在文本块上方且距离较近
            if abs(img.get('y1', 0) - text_y) < threshold:
                nearby.append(img)
                img['processed'] = True

        return nearby

    def _image_to_placeholder(self, image: Dict) -> str:
        """将图片转换为占位符"""
        # 根据图片大小和位置生成描述
        width = image.get('x1', 0) - image.get('x0', 0)
        height = image.get('y1', 0) - image.get('y0', 0)

        # 简单的大小分类
        if width > 400 and height > 300:
            size_desc = "大图"
        elif width > 200 and height > 150:
            size_desc = "中等图片"
        else:
            size_desc = "小图"

        return f"\n[图片：文档中的{size_desc}]\n"

    def _clean_content(self, content: str) -> str:
        """清理内容格式"""
        # 移除多余的空行
        content = re.sub(r'\n{3,}', '\n\n', content)
        # 移除行尾空格
        content = re.sub(r' +\n', '\n', content)
        return content.strip()


class WordExtractor(BaseExtractor):
    """Word文档提取器 - 使用python-docx提取结构化内容"""

    def __init__(self):
        self._Document = None

    def _get_document_class(self):
        """延迟加载 docx 模块"""
        if self._Document is None:
            try:
                from docx import Document
                self._Document = Document
            except ImportError:
                raise ImportError("python-docx 未安装，请运行: pip install python-docx")
        return self._Document

    def extract(self, file_path: Path) -> ExtractedContent:
        """提取Word文档的结构化内容"""
        Document = self._get_document_class()
        doc = Document(str(file_path))

        try:
            content_parts = []
            current_list_type = None  # 'bullet', 'numbered', None

            # 统计图片数量（用于元数据）
            image_count = len(doc.inline_shapes)

            for para in doc.paragraphs:
                if not para.text.strip():
                    # 空段落表示列表结束
                    if current_list_type:
                        current_list_type = None
                    continue

                md_text = self._paragraph_to_markdown(para)

                if md_text:
                    content_parts.append(md_text)

            # 处理表格
            for i, table in enumerate(doc.tables, 1):
                table_md = self._convert_table_to_markdown(table)
                if table_md:
                    content_parts.append(f"\n### 表格 {i}\n")
                    content_parts.append(table_md)

            # 添加图片占位符（如果文档中有图片但未在表格中）
            if image_count > 0:
                content_parts.append(f"\n<!-- 文档包含 {image_count} 张嵌入图片 -->\n")

            content = "\n\n".join(content_parts)
            content = self._clean_content(content)

            # 元数据
            core_props = doc.core_properties
            metadata = {
                '标题': core_props.title,
                '作者': core_props.author,
                '创建时间': core_props.created.strftime('%Y-%m-%d') if core_props.created else None,
                '修改时间': core_props.modified.strftime('%Y-%m-%d') if core_props.modified else None,
                '段落数': len(doc.paragraphs),
                '表格数': len(doc.tables),
                '图片数': image_count,
            }

            header = self._create_header(file_path, "Word文档", metadata)
            full_content = header + "\n\n" + content

            return ExtractedContent(
                title=core_props.title or file_path.stem,
                content=full_content,
                metadata=metadata,
                source_type='word'
            )

        finally:
            pass

    def _paragraph_to_markdown(self, para) -> str:
        """将段落转换为Markdown"""
        text = para.text.strip()
        if not text:
            return ""

        style_name = para.style.name if para.style else ""

        # 标题样式
        if style_name.startswith('Heading 1') or style_name.startswith('标题 1'):
            return f"# {text}"
        elif style_name.startswith('Heading 2') or style_name.startswith('标题 2'):
            return f"## {text}"
        elif style_name.startswith('Heading 3') or style_name.startswith('标题 3'):
            return f"### {text}"
        elif style_name.startswith('Heading 4') or style_name.startswith('标题 4'):
            return f"#### {text}"

        # 列表样式
        if style_name.startswith('List Bullet') or style_name.startswith('项目符号'):
            return f"- {text}"
        if style_name.startswith('List Number') or style_name.startswith('编号'):
            # 移除原有的数字
            text = re.sub(r'^\d+[\.\)]\s*', '', text)
            return f"1. {text}"

        # 检查是否为编号段落（如 "1. 内容" 或 "(1) 内容"）
        if re.match(r'^\d+[\.\)]\s+', text):
            return text

        # 检查是否为项目符号段落
        if text.startswith(('• ', '· ', '● ')):
            return f"- {text[2:]}"

        # 普通段落
        # 处理段落中的粗体和斜体
        return self._process_inline_formatting(para)

    def _process_inline_formatting(self, para) -> str:
        """处理段落内的格式（粗体、斜体）"""
        result = []

        for run in para.runs:
            text = run.text
            if not text:
                continue

            # 处理粗体和斜体
            if run.bold and run.italic:
                result.append(f"***{text}***")
            elif run.bold:
                result.append(f"**{text}**")
            elif run.italic:
                result.append(f"*{text}*")
            else:
                result.append(text)

        return "".join(result)

    def _convert_table_to_markdown(self, table) -> str:
        """将Word表格转换为Markdown"""
        if not table.rows:
            return ""

        lines = []

        for i, row in enumerate(table.rows):
            cells = [self._escape_table_cell(cell.text.strip()) for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

            # 在第一行后添加分隔行
            if i == 0:
                separator = "|" + "---|" * len(row.cells)
                lines.append(separator)

        return "\n".join(lines)

    def _clean_content(self, content: str) -> str:
        """清理内容格式"""
        content = re.sub(r'\n{3,}', '\n\n', content)
        content = re.sub(r' +\n', '\n', content)
        return content.strip()


class PPTXExtractor(BaseExtractor):
    """PowerPoint提取器 - 使用python-pptx提取幻灯片内容"""

    def __init__(self):
        self._Presentation = None

    def _get_presentation_class(self):
        """延迟加载 pptx 模块"""
        if self._Presentation is None:
            try:
                from pptx import Presentation
                self._Presentation = Presentation
            except ImportError:
                raise ImportError("python-pptx 未安装，请运行: pip install python-pptx")
        return self._Presentation

    def extract(self, file_path: Path) -> ExtractedContent:
        """提取PPTX的结构化内容"""
        Presentation = self._get_presentation_class()
        prs = Presentation(str(file_path))

        content_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = self._extract_slide_content(slide, slide_num)
            if slide_content:
                content_parts.append(slide_content)

        content = "\n\n".join(content_parts)
        content = self._clean_content(content)

        # 元数据
        core_props = prs.core_properties
        metadata = {
            '标题': core_props.title,
            '作者': core_props.author,
            '幻灯片数': len(prs.slides),
        }

        header = self._create_header(file_path, "PowerPoint", metadata)
        full_content = header + "\n\n" + content

        return ExtractedContent(
            title=core_props.title or file_path.stem,
            content=full_content,
            metadata=metadata,
            source_type='pptx'
        )

    def _extract_slide_content(self, slide, slide_num: int) -> str:
        """提取单张幻灯片的内容"""
        parts = []

        # 幻灯片标题
        title = self._get_slide_title(slide)
        if title:
            parts.append(f"## 幻灯片 {slide_num}：{title}")
        else:
            parts.append(f"## 幻灯片 {slide_num}")

        parts.append("")

        # 提取所有形状的内容
        for shape in slide.shapes:
            shape_content = self._extract_shape_content(shape)
            if shape_content:
                parts.append(shape_content)

        return "\n".join(parts)

    def _get_slide_title(self, slide) -> str:
        """获取幻灯片标题"""
        if slide.shapes.title:
            return slide.shapes.title.text.strip()

        # 尝试从第一个文本形状获取
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                return shape.text.strip()

        return ""

    def _extract_shape_content(self, shape) -> str:
        """提取形状的内容"""
        # 检查是否为图片 (MSO_SHAPE_TYPE.PICTURE = 13)
        if shape.shape_type == 13:
            # 尝试获取图片描述（替代文本）
            description = ""
            if hasattr(shape, 'image'):
                description = getattr(shape, 'name', '图片')
            return f"\n[图片：{description or '幻灯片中的图片'}]\n"

        # 检查是否为图表 (MSO_SHAPE_TYPE.CHART = 3)
        if shape.shape_type == 3 and hasattr(shape, 'chart'):
            return self._extract_chart_data(shape.chart, shape.name)

        # 检查是否为SmartArt (MSO_SHAPE_TYPE.SMARTART = 6)
        if shape.shape_type == 6:
            return self._extract_smartart_data(shape)

        # 文本框
        if hasattr(shape, "text") and shape.text.strip():
            text = shape.text.strip()

            # 检查是否为表格
            if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
                return self._convert_pptx_table_to_markdown(shape.table)

            # 普通文本
            return self._format_slide_text(text)

        # 表格
        if shape.has_table:
            return self._convert_pptx_table_to_markdown(shape.table)

        return ""

    def _extract_chart_data(self, chart, shape_name: str = "图表") -> str:
        """提取图表数据并转换为Markdown表格"""
        lines = [f"\n### 图表：{shape_name}\n"]

        try:
            # 获取图表标题
            if hasattr(chart, 'chart_title') and chart.chart_title:
                title_text = chart.chart_title.text_frame.text if chart.chart_title.has_text_frame else ""
                if title_text:
                    lines.append(f"**{title_text}**\n")

            # 获取图表数据
            if hasattr(chart, 'plot') and chart.plot:
                plot = chart.plot

                # 处理不同图表类型
                if hasattr(plot, 'series'):
                    # 收集所有数据
                    categories = []
                    series_data = {}

                    for series in plot.series:
                        series_name = series.name if series.name else "系列"
                        series_data[series_name] = []

                        # 获取类别（X轴标签）
                        if hasattr(series, 'categories'):
                            for cat in series.categories:
                                cat_label = str(cat.label) if hasattr(cat, 'label') else str(cat)
                                if cat_label not in categories:
                                    categories.append(cat_label)

                        # 获取数值
                        if hasattr(series, 'values'):
                            for val in series.values:
                                series_data[series_name].append(str(val) if val is not None else "")

                    # 构建Markdown表格
                    if categories and series_data:
                        # 表头
                        header_cells = ["类别"] + list(series_data.keys())
                        lines.append("| " + " | ".join(header_cells) + " |")
                        lines.append("|" + "---|" * len(header_cells))

                        # 数据行
                        for i, cat in enumerate(categories):
                            row = [str(cat)]
                            for series_name in series_data.keys():
                                values = series_data[series_name]
                                val = values[i] if i < len(values) else ""
                                row.append(val)
                            lines.append("| " + " | ".join(row) + " |")
                    elif series_data:
                        # 没有类别，只有数值
                        for series_name, values in series_data.items():
                            lines.append(f"- **{series_name}**: {', '.join(values)}")
                else:
                    lines.append("*[图表数据结构不支持自动提取]*")

        except Exception as e:
            lines.append(f"*[图表数据提取失败: {str(e)}]*")

        return "\n".join(lines) + "\n"

    def _extract_smartart_data(self, shape) -> str:
        """提取SmartArt的层级结构和文本"""
        lines = [f"\n### 结构图：{getattr(shape, 'name', 'SmartArt')}\n"]

        try:
            # SmartArt在python-pptx中的访问有限
            # 尝试提取所有文本节点
            texts = []

            def extract_texts_from_shape(shp, level=0):
                """递归提取形状中的文本"""
                if hasattr(shp, 'text') and shp.text.strip():
                    texts.append((level, shp.text.strip()))

                # 如果是组合形状，递归处理
                if hasattr(shp, 'shapes'):
                    for child in shp.shapes:
                        extract_texts_from_shape(child, level + 1)

            extract_texts_from_shape(shape)

            if texts:
                # 构建层级结构
                for level, text in texts:
                    indent = "  " * level
                    lines.append(f"{indent}- {text}")
            else:
                lines.append("*[SmartArt结构图]*")

        except Exception as e:
            lines.append(f"*[SmartArt提取失败: {str(e)}]*")

        return "\n".join(lines) + "\n"

    def _is_chart_or_smartart(self, shape) -> bool:
        """检查是否为图表或SmartArt（MSO_SHAPE_TYPE）"""
        # 3 = CHART, 6 = SMARTART
        return shape.shape_type in [3, 6]

    def _get_shape_type_name(self, shape_type: int) -> str:
        """获取形状类型名称"""
        type_names = {
            3: "图表",
            6: "SmartArt图形",
            13: "图片",
            17: "媒体",
            20: "OLE对象",
        }
        return type_names.get(shape_type, "图形")

    def _format_slide_text(self, text: str) -> str:
        """格式化幻灯片文本"""
        lines = text.split('\n')
        formatted_lines = []

        for line in lines:
            line = line.strip()
            if not line:
                continue

            # 检测项目符号
            if line.startswith(('•', '·', '●', '-', '*')):
                line = f"- {line[1:].strip()}"
            # 检测编号
            elif re.match(r'^\d+[\.\)]\s', line):
                pass  # 保留原样

            formatted_lines.append(line)

        return "\n".join(formatted_lines)

    def _convert_pptx_table_to_markdown(self, table) -> str:
        """将PPTX表格转换为Markdown"""
        if not table.rows:
            return ""

        lines = []

        for i, row in enumerate(table.rows):
            cells = [self._escape_table_cell(cell.text.strip()) for cell in row.cells]
            lines.append("| " + " | ".join(cells) + " |")

            # 在第一行后添加分隔行
            if i == 0:
                separator = "|" + "---|" * len(row.cells)
                lines.append(separator)

        return "\n".join(lines)

    def _clean_content(self, content: str) -> str:
        """清理内容格式"""
        content = re.sub(r'\n{4,}', '\n\n\n', content)
        content = re.sub(r' +\n', '\n', content)
        return content.strip()


class ExtractorFactory:
    """提取器工厂"""

    _extractors = {
        'vector_pdf': VectorPDFExtractor,
        'word': WordExtractor,
        'pptx': PPTXExtractor,
    }

    @classmethod
    def get_extractor(cls, doc_type: str) -> BaseExtractor:
        """获取对应类型的提取器"""
        extractor_class = cls._extractors.get(doc_type)
        if extractor_class:
            return extractor_class()
        raise ValueError(f"不支持的文档类型: {doc_type}")

    @classmethod
    def extract_file(cls, file_path: Path, doc_type: str) -> ExtractedContent:
        """便捷方法：提取指定文件"""
        extractor = cls.get_extractor(doc_type)
        return extractor.extract(file_path)
