"""
文档类型检测模块

用于检测文件类型，区分：
- 扫描件PDF（图片型，需要OCR）
- 文本型PDF（可直接提取文字）
- 其他格式（PPTX、Word、MD等）
"""

from enum import Enum
from dataclasses import dataclass
from pathlib import Path
from typing import List, Optional, Tuple


class DocumentType(Enum):
    """文档类型枚举"""
    SCANNED_PDF = "scanned_pdf"      # 扫描件PDF，需要OCR
    VECTOR_PDF = "vector_pdf"        # 矢量PDF，可提取结构化信息
    WORD = "word"                    # Word文档 (docx)
    PPTX = "pptx"                    # PowerPoint演示文稿
    UNKNOWN = "unknown"              # 未知类型


@dataclass
class DetectionResult:
    """文档检测结果"""
    file_path: Path
    doc_type: DocumentType
    reason: str                      # 判断原因说明
    confidence: float                # 置信度 0-1
    page_count: int = 0              # 页数（PDF/PPTX）
    avg_chars_per_page: float = 0.0  # 平均每页字符数
    text_page_ratio: float = 0.0     # 有文本的页面占比
    has_structured_content: bool = False  # 是否有结构化内容（表格、标题等）


class DocumentDetector:
    """文档类型检测器"""

    # 支持的扩展名映射
    EXTENSION_MAP = {
        '.pdf': 'pdf',
        '.docx': 'word',
        '.doc': 'word',
        '.pptx': 'pptx',
        '.ppt': 'pptx',
    }

    # 扫描件判断阈值
    MIN_CHARS_PER_PAGE = 100         # 平均每页最少字符数（低于此值可能是扫描件）
    MIN_TEXT_PAGE_RATIO = 0.3        # 最少文本页占比（低于此值可能是扫描件）
    TEXT_PAGE_MIN_CHARS = 500        # 单页认定为文本页的最少字符数
    VECTOR_PDF_MIN_CHARS = 200       # 矢量PDF最少字符数阈值

    def __init__(self):
        self._fitz = None

    def _get_fitz(self):
        """延迟加载 fitz 模块"""
        if self._fitz is None:
            try:
                import fitz
                self._fitz = fitz
            except ImportError:
                raise ImportError(
                    "pymupdf 未安装，请运行: pip install pymupdf"
                )
        return self._fitz

    def detect(self, file_path: Path) -> DetectionResult:
        """
        检测单个文件的类型

        Args:
            file_path: 文件路径

        Returns:
            DetectionResult 检测结果
        """
        file_path = Path(file_path)

        if not file_path.exists():
            return DetectionResult(
                file_path=file_path,
                doc_type=DocumentType.UNKNOWN,
                reason="文件不存在",
                confidence=1.0
            )

        # 检查扩展名
        ext = file_path.suffix.lower()
        doc_category = self.EXTENSION_MAP.get(ext)

        if not doc_category:
            return DetectionResult(
                file_path=file_path,
                doc_type=DocumentType.UNKNOWN,
                reason=f"不支持的文件类型 ({ext})",
                confidence=1.0
            )

        # 根据类型分发检测
        if doc_category == 'pdf':
            return self._detect_pdf_type(file_path)
        elif doc_category == 'word':
            return self._detect_word_type(file_path)
        elif doc_category == 'pptx':
            return self._detect_pptx_type(file_path)

        return DetectionResult(
            file_path=file_path,
            doc_type=DocumentType.UNKNOWN,
            reason=f"无法识别的文档类型 ({ext})",
            confidence=0.0
        )

    def _detect_pdf_type(self, pdf_path: Path) -> DetectionResult:
        """
        检测PDF是扫描件还是文本型

        判断逻辑：
        1. 无文本层 + 有图片 → 扫描件
        2. 平均每页字符数 < 100 → 扫描件
        3. 文本页占比 < 30% → 扫描件
        4. 其他 → 文本型PDF
        """
        fitz = self._get_fitz()

        try:
            doc = fitz.open(str(pdf_path))
            total_pages = len(doc)

            if total_pages == 0:
                doc.close()
                return DetectionResult(
                    file_path=pdf_path,
                    doc_type=DocumentType.UNKNOWN,
                    reason="PDF为空",
                    confidence=1.0,
                    page_count=0
                )

            text_pages = 0
            image_pages = 0
            total_chars = 0

            for page in doc:
                # 获取文本
                text = page.get_text()
                char_count = len(text.strip())
                total_chars += char_count

                # 判断是否为文本页
                if char_count >= self.TEXT_PAGE_MIN_CHARS:
                    text_pages += 1

                # 检查是否包含图片
                images = page.get_images()
                if len(images) > 0:
                    image_pages += 1

            doc.close()

            # 计算指标
            avg_chars = total_chars / total_pages if total_pages > 0 else 0
            text_ratio = text_pages / total_pages if total_pages > 0 else 0

            # 判断是否为扫描件
            is_scanned = False
            reasons = []

            # 规则1: 无文本层但有图片
            if text_pages == 0 and image_pages > 0:
                is_scanned = True
                reasons.append(f"无文本层，{image_pages}页含图片")

            # 规则2: 平均每页字符数过少
            if avg_chars < self.MIN_CHARS_PER_PAGE:
                is_scanned = True
                reasons.append(f"平均每页仅{avg_chars:.0f}字符")

            # 规则3: 文本页占比过低
            if text_ratio < self.MIN_TEXT_PAGE_RATIO:
                is_scanned = True
                reasons.append(f"仅{text_ratio*100:.0f}%页面有文本")

            # 计算置信度
            if is_scanned:
                confidence = min(0.95, 0.7 + len(reasons) * 0.1)
                doc_type = DocumentType.SCANNED_PDF
                reason = "扫描件: " + "; ".join(reasons)
                has_structured = False
            else:
                confidence = min(0.95, 0.7 + text_ratio * 0.2)
                doc_type = DocumentType.VECTOR_PDF
                reason = f"矢量PDF: 平均每页{avg_chars:.0f}字符，{text_ratio*100:.0f}%页面有文本"
                has_structured = avg_chars >= self.VECTOR_PDF_MIN_CHARS

            return DetectionResult(
                file_path=pdf_path,
                doc_type=doc_type,
                reason=reason,
                confidence=confidence,
                page_count=total_pages,
                avg_chars_per_page=avg_chars,
                text_page_ratio=text_ratio,
                has_structured_content=has_structured
            )

        except Exception as e:
            return DetectionResult(
                file_path=pdf_path,
                doc_type=DocumentType.UNKNOWN,
                reason=f"检测失败: {str(e)}",
                confidence=0.0
            )

    def _detect_word_type(self, doc_path: Path) -> DetectionResult:
        """
        检测Word文档类型

        尝试读取文档获取基本信息
        """
        try:
            from docx import Document

            doc = Document(str(doc_path))

            # 统计段落数
            para_count = len(doc.paragraphs)
            table_count = len(doc.tables)

            # 尝试获取文档属性
            core_props = doc.core_properties
            title = core_props.title or ""

            doc.close()

            has_structure = table_count > 0 or para_count > 50

            return DetectionResult(
                file_path=doc_path,
                doc_type=DocumentType.WORD,
                reason=f"Word文档: {para_count}段落, {table_count}表格",
                confidence=1.0,
                page_count=0,  # Word页数难以准确获取
                has_structured_content=has_structure
            )

        except Exception as e:
            return DetectionResult(
                file_path=doc_path,
                doc_type=DocumentType.WORD,
                reason=f"Word文档 (无法读取详情: {str(e)})",
                confidence=0.9
            )

    def _detect_pptx_type(self, pptx_path: Path) -> DetectionResult:
        """
        检测PPTX文档类型

        尝试读取演示文稿获取基本信息
        """
        try:
            from pptx import Presentation

            prs = Presentation(str(pptx_path))
            slide_count = len(prs.slides)

            # 统计总形状数（文本框、表格等）
            total_shapes = 0
            for slide in prs.slides:
                total_shapes += len(slide.shapes)

            has_structure = total_shapes > 0

            return DetectionResult(
                file_path=pptx_path,
                doc_type=DocumentType.PPTX,
                reason=f"PowerPoint演示文稿: {slide_count}张幻灯片, {total_shapes}个元素",
                confidence=1.0,
                page_count=slide_count,
                has_structured_content=has_structure
            )

        except Exception as e:
            return DetectionResult(
                file_path=pptx_path,
                doc_type=DocumentType.PPTX,
                reason=f"PowerPoint演示文稿 (无法读取详情: {str(e)})",
                confidence=0.9
            )

    def detect_batch(self, directory: Path, recursive: bool = False) -> List[DetectionResult]:
        """
        批量检测目录下的所有文件

        Args:
            directory: 目录路径
            recursive: 是否递归子目录

        Returns:
            DetectionResult 列表
        """
        directory = Path(directory)
        results = []

        if recursive:
            pattern = "**/*"
        else:
            pattern = "*"

        for file_path in directory.glob(pattern):
            if file_path.is_file():
                result = self.detect(file_path)
                results.append(result)

        return results

    def filter_scanned_pdfs(self, directory: Path, recursive: bool = False) -> Tuple[List[Path], List[DetectionResult]]:
        """
        筛选出需要处理的扫描件PDF

        Args:
            directory: 目录路径
            recursive: 是否递归子目录

        Returns:
            (scanned_pdf_paths, all_results)
        """
        results = self.detect_batch(directory, recursive)

        scanned_pdfs = [
            r.file_path for r in results
            if r.doc_type == DocumentType.SCANNED_PDF
        ]

        return scanned_pdfs, results

    def print_detection_report(self, results: List[DetectionResult]):
        """打印检测结果报告"""
        # 分类统计
        scanned = [r for r in results if r.doc_type == DocumentType.SCANNED_PDF]
        vector_pdfs = [r for r in results if r.doc_type == DocumentType.VECTOR_PDF]
        words = [r for r in results if r.doc_type == DocumentType.WORD]
        pptxs = [r for r in results if r.doc_type == DocumentType.PPTX]
        unknown = [r for r in results if r.doc_type == DocumentType.UNKNOWN]

        print("\n" + "=" * 60)
        print("批量处理分析:")
        print("=" * 60)

        # 扫描件（需要OCR）
        if scanned:
            print(f"\n📄 扫描件PDF (Vision API识别): {len(scanned)} 个")
            for r in scanned:
                print(f"   ✓ {r.file_path.name}")

        # 矢量PDF（Vision API识别）
        if vector_pdfs:
            print(f"\n📝 矢量PDF (Vision API识别): {len(vector_pdfs)} 个")
            for r in vector_pdfs:
                print(f"   ✓ {r.file_path.name}")

        # Word文档（结构化提取）
        if words:
            print(f"\n📘 Word文档 (结构化提取): {len(words)} 个")
            for r in words:
                print(f"   ✓ {r.file_path.name}")

        # PPTX（Vision API识别）
        if pptxs:
            print(f"\n📊 PowerPoint (Vision API识别): {len(pptxs)} 个")
            for r in pptxs:
                print(f"   ✓ {r.file_path.name}")

        # 未知/错误
        if unknown:
            print(f"\n❓ 未知/错误: {len(unknown)} 个")
            for r in unknown:
                print(f"   ? {r.file_path.name} ({r.reason})")

        print("\n" + "=" * 60)


def extract_text_from_pdf(pdf_path: Path, output_path: Optional[Path] = None) -> str:
    """
    从文本型PDF中提取文字

    Args:
        pdf_path: PDF文件路径
        output_path: 输出文件路径（可选）

    Returns:
        提取的文本内容
    """
    try:
        import fitz
    except ImportError:
        raise ImportError("pymupdf 未安装，请运行: pip install pymupdf")

    doc = fitz.open(str(pdf_path))
    text_parts = []

    for page_num, page in enumerate(doc, 1):
        text = page.get_text()
        if text.strip():
            text_parts.append(f"\n--- 第{page_num}页 ---\n")
            text_parts.append(text)

    doc.close()

    full_text = "".join(text_parts)

    # 保存到文件
    if output_path:
        output_path = Path(output_path)
        output_path.write_text(full_text, encoding='utf-8')

    return full_text
